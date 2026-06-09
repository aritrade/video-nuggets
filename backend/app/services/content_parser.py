"""
Content parser supporting PDF, PPTX, TXT, PNG/JPG (OCR), and URL extraction.
Extracts structured text content from various source formats.
"""
import re
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional

import httpx
from bs4 import BeautifulSoup
from PIL import Image
import pytesseract
from PyPDF2 import PdfReader
from pptx import Presentation


@dataclass
class ContentSection:
    title: str
    body: str
    key: str = ""
    subsections: list["ContentSection"] = field(default_factory=list)
    images: list[str] = field(default_factory=list)
    # Optional pin: force this section to use a specific Cloud Bible diagram
    # (manifest id, e.g. "p007_00") and/or a specific slide layout.
    diagram_id: str = ""
    preferred_layout: str = ""


@dataclass
class ParsedContent:
    title: str
    sections: list[ContentSection]
    source_type: str
    raw_text: str = ""


def parse_source(source_path: str) -> ParsedContent:
    """Route to the appropriate parser based on source type."""
    if source_path.startswith("http://") or source_path.startswith("https://"):
        return parse_url(source_path)

    path = Path(source_path)
    ext = path.suffix.lower()
    parsers = {
        ".pdf": parse_pdf,
        ".pptx": parse_pptx,
        ".txt": parse_txt,
        ".png": parse_image,
        ".jpg": parse_image,
        ".jpeg": parse_image,
    }
    parser = parsers.get(ext)
    if not parser:
        raise ValueError(f"Unsupported file format: {ext}")
    return parser(source_path)


def parse_pdf(file_path: str) -> ParsedContent:
    reader = PdfReader(file_path)
    sections = []
    current_title = ""
    current_body = []
    full_text = []

    for page in reader.pages:
        text = page.extract_text() or ""
        full_text.append(text)
        lines = text.split("\n")

        for line in lines:
            line = line.strip()
            if not line:
                continue
            if _is_heading(line):
                if current_title or current_body:
                    sections.append(ContentSection(
                        title=current_title or "Introduction",
                        body="\n".join(current_body),
                        key=_make_key(current_title),
                    ))
                current_title = line
                current_body = []
            else:
                current_body.append(line)

    if current_title or current_body:
        sections.append(ContentSection(
            title=current_title or "Content",
            body="\n".join(current_body),
            key=_make_key(current_title),
        ))

    if not sections:
        sections = [ContentSection(
            title="Document Content",
            body="\n".join(full_text),
            key="document_content",
        )]

    sections = _consolidate_sections(sections)

    return ParsedContent(
        title=sections[0].title if sections else "Untitled",
        sections=sections,
        source_type="pdf",
        raw_text="\n".join(full_text),
    )


def _consolidate_sections(
    sections: list[ContentSection],
    min_body_chars: int = 400,
    max_sections: int = 40,
) -> list[ContentSection]:
    """
    Merge tiny adjacent sections so the LLM simplifier doesn't get hammered with
    hundreds of trivial chunks. Two-pass:
      1. Roll any section whose body is below `min_body_chars` into the next one
         (keeps the small section's title inline as a sub-heading).
      2. If we still exceed `max_sections`, repeatedly merge the shortest pair
         of adjacent sections until we're under the cap.
    """
    if not sections:
        return sections

    merged: list[ContentSection] = []
    pending: ContentSection | None = None

    for sec in sections:
        if pending is not None:
            sec = ContentSection(
                title=pending.title or sec.title,
                body=(pending.body + ("\n\n" + sec.title + "\n" if sec.title else "\n") + sec.body).strip(),
                key=pending.key or sec.key,
                subsections=pending.subsections + sec.subsections,
                images=pending.images + sec.images,
            )
            pending = None

        if len(sec.body.strip()) < min_body_chars:
            pending = sec
            continue

        merged.append(sec)

    if pending is not None:
        if merged:
            tail = merged[-1]
            merged[-1] = ContentSection(
                title=tail.title,
                body=(tail.body + "\n\n" + (pending.title + "\n" if pending.title else "") + pending.body).strip(),
                key=tail.key,
                subsections=tail.subsections + pending.subsections,
                images=tail.images + pending.images,
            )
        else:
            merged.append(pending)

    while len(merged) > max_sections:
        shortest_idx = min(
            range(len(merged) - 1),
            key=lambda i: len(merged[i].body) + len(merged[i + 1].body),
        )
        a, b = merged[shortest_idx], merged[shortest_idx + 1]
        combined = ContentSection(
            title=a.title,
            body=(a.body + "\n\n" + (b.title + "\n" if b.title else "") + b.body).strip(),
            key=a.key,
            subsections=a.subsections + b.subsections,
            images=a.images + b.images,
        )
        merged[shortest_idx:shortest_idx + 2] = [combined]

    return merged


def parse_pptx(file_path: str) -> ParsedContent:
    prs = Presentation(file_path)
    sections = []
    full_text = []

    for slide in prs.slides:
        slide_title = ""
        slide_body = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    if not slide_title:
                        slide_title = text
                    else:
                        slide_body.append(text)
                    full_text.append(text)

        if slide_title:
            sections.append(ContentSection(
                title=slide_title,
                body="\n".join(slide_body),
                key=_make_key(slide_title),
            ))

    return ParsedContent(
        title=sections[0].title if sections else "Presentation",
        sections=sections,
        source_type="pptx",
        raw_text="\n".join(full_text),
    )


_MIN_BODY_CHARS = 220
_MAX_SECTIONS = 8


def _txt_is_heading(line: str) -> bool:
    """Detect a section heading: a short, title-cased line without sentence-
    terminal punctuation (a trailing '?' is allowed, since questions are common
    titles), or an explicit markdown/numbered heading."""
    l = line.strip()
    if not (4 <= len(l) <= 80):
        return False
    if re.match(r"^(#{1,3}\s+|chapter\s+\d+|section\s+\d+|part\s+[ivx0-9]+)\b", l, re.IGNORECASE):
        return True
    if re.match(r"^\d+(\.\d+){0,2}\s+[A-Z][A-Za-z]", l):
        return True
    ends_clean = not re.search(r"[.,;:!]$", l)
    words = l.split()
    if ends_clean and 2 <= len(words) <= 12:
        caps = sum(1 for w in words if w[:1].isupper())
        if caps / len(words) >= 0.5:
            return True
    return False


def _txt_consolidate(sections: list[ContentSection]) -> list[ContentSection]:
    """Merge too-short sections into a neighbour, then cap the section count."""
    if not sections:
        return sections
    merged: list[ContentSection] = []
    for sec in sections:
        if merged and len(sec.body.strip()) < _MIN_BODY_CHARS:
            prev = merged[-1]
            extra = (("\n" + sec.title + "\n") if sec.title else "\n") + sec.body
            prev.body = (prev.body + extra).strip()
        else:
            merged.append(sec)
    while len(merged) > _MAX_SECTIONS:
        # Merge the pair of adjacent sections with the least combined text.
        best_i, best_len = 0, 10 ** 9
        for i in range(len(merged) - 1):
            combined = len(merged[i].body) + len(merged[i + 1].body)
            if combined < best_len:
                best_len, best_i = combined, i
        a, b = merged[best_i], merged[best_i + 1]
        a.body = (a.body + "\n\n" + (b.title + "\n" if b.title else "") + b.body).strip()
        del merged[best_i + 1]
    return merged


def _parse_text_blocks(text: str, source_type: str) -> ParsedContent:
    lines = text.replace("\r", "").split("\n")
    sections: list[ContentSection] = []
    state = {"title": "", "body": []}

    def flush():
        if state["title"] or state["body"]:
            title = state["title"] or "Introduction"
            sections.append(ContentSection(
                title=title,
                body="\n".join(state["body"]).strip(),
                key=_make_key(title),
            ))

    for raw in lines:
        line = raw.strip()
        if not line:
            continue
        if _txt_is_heading(line):
            flush()
            state = {"title": re.sub(r"^#{1,3}\s+", "", line), "body": []}
        else:
            state["body"].append(line)
    flush()

    sections = _txt_consolidate(sections)
    title = sections[0].title if sections else "Text Document"
    return ParsedContent(title=title, sections=sections, source_type=source_type, raw_text=text)


def parse_txt(file_path: str) -> ParsedContent:
    text = Path(file_path).read_text(encoding="utf-8", errors="ignore")
    return _parse_text_blocks(text, source_type="txt")


def parse_image(file_path: str) -> ParsedContent:
    image = Image.open(file_path)
    text = pytesseract.image_to_string(image)

    sections = [ContentSection(
        title="Extracted from Image",
        body=text.strip(),
        key="image_content",
        images=[file_path],
    )]

    return ParsedContent(
        title="Image Content",
        sections=sections,
        source_type="image",
        raw_text=text,
    )


def parse_url(url: str) -> ParsedContent:
    response = httpx.get(url, follow_redirects=True, timeout=30.0)
    response.raise_for_status()

    soup = BeautifulSoup(response.text, "html.parser")

    for tag in soup(["script", "style", "nav", "footer", "header"]):
        tag.decompose()

    title = soup.title.string if soup.title else "Web Content"
    sections = []

    headings = soup.find_all(["h1", "h2", "h3"])
    if headings:
        for heading in headings:
            body_parts = []
            sibling = heading.find_next_sibling()
            while sibling and sibling.name not in ["h1", "h2", "h3"]:
                text = sibling.get_text(strip=True)
                if text:
                    body_parts.append(text)
                sibling = sibling.find_next_sibling()
            sections.append(ContentSection(
                title=heading.get_text(strip=True),
                body="\n".join(body_parts),
                key=_make_key(heading.get_text(strip=True)),
            ))
    else:
        body_text = soup.get_text(separator="\n", strip=True)
        sections = [ContentSection(
            title=title,
            body=body_text,
            key="web_content",
        )]

    raw_text = soup.get_text(separator="\n", strip=True)

    return ParsedContent(
        title=title,
        sections=sections,
        source_type="url",
        raw_text=raw_text,
    )


def _is_heading(line: str) -> bool:
    if len(line) > 80 or len(line) < 4:
        return False
    if re.match(r"^(#{1,3}\s+|Chapter\s+\d+|Section\s+\d+|Part\s+[IVX0-9]+|Appendix\s+[A-Z0-9])\b", line):
        return True
    if re.match(r"^\d+(\.\d+){0,2}\s+[A-Z][A-Za-z]", line):
        return True
    if line.isupper() and len(line.split()) >= 2 and not line.endswith((".", ",", ";")):
        return True
    return False


def _make_key(title: str) -> str:
    key = re.sub(r"[^a-z0-9]+", "_", title.lower().strip())
    return key[:80].strip("_")
