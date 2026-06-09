"""
Neutral slide-deck generator using python-pptx's built-in template.

This produces the optional *downloadable* PPTX deck only. The actual video
frames are rendered separately by `slide_image_generator` (PIL), so this deck
is purely a convenience export and carries no proprietary template.
"""
from pptx import Presentation
from pptx.util import Pt

from app.config import SLIDES_DIR
from app.services.content_parser import ParsedContent, ContentSection

# python-pptx default template layout indices.
_LAYOUT_TITLE = 0          # Title + subtitle
_LAYOUT_TITLE_CONTENT = 1  # Title + bulleted content
_LAYOUT_SECTION = 2        # Section header
_LAYOUT_BLANK = 6          # Blank


def generate_slides(content: ParsedContent, video_id: int, visualizations: list[str] = None) -> str:
    """Generate a neutral PPTX from parsed/simplified content."""
    prs = Presentation()  # default blank template, no proprietary assets

    _add_title_slide(prs, content.title)

    for i, section in enumerate(content.sections):
        _add_content_slide(prs, section)

    _add_closing_slide(prs)

    output_path = str(SLIDES_DIR / f"video_{video_id}.pptx")
    prs.save(output_path)
    return output_path


def _add_title_slide(prs: Presentation, title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE])
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "A Video Nugget"


def _add_content_slide(prs: Presentation, section: ContentSection):
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_TITLE_CONTENT])
    slide.shapes.title.text = section.title
    body = slide.placeholders[1].text_frame
    body.clear()
    lines = [l.strip() for l in section.body.split("\n") if l.strip()]
    for idx, line in enumerate(lines[:8]):
        para = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        para.text = line
        para.space_after = Pt(6)


def _add_closing_slide(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[_LAYOUT_SECTION])
    slide.shapes.title.text = "Thanks for watching!"
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = "Made with Video Nuggets OS"
